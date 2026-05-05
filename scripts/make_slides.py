"""
Generate PowerPoint slide decks for all three MLS NLP papers.
Style matches the V2 NBA Presentation: white backgrounds, dark text,
card layouts, numbered findings, footer on every slide.

Output: presentations/paper1_slides.pptx
        presentations/paper2_slides.pptx
        presentations/research_note_slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT_DIR = Path("presentations")
OUT_DIR.mkdir(exist_ok=True)

# ── Color palette (matches NBA deck) ─────────────────────────────────────────
DARK_NAVY  = RGBColor(0x0D, 0x1A, 0x33)   # titles
MID_BLUE   = RGBColor(0x1D, 0x42, 0x8A)   # section headers / card headers
GRAY_TEXT  = RGBColor(0x44, 0x55, 0x77)   # body text
LIGHT_GRAY = RGBColor(0xBB, 0xCC, 0xFF)   # light accent text
RED        = RGBColor(0xC8, 0x10, 0x2E)   # highlight numbers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF7, 0xF8, 0xFC)   # card bg
CARD_DARK  = RGBColor(0x0D, 0x1A, 0x33)   # dark stat card bg
FOOTER_BG  = RGBColor(0x0D, 0x1A, 0x33)

FOOTER_TEXT = "MLS NLP Project  ·  Dante  ·  Wilson  ·  Faisal  ·  Fidel"

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=WHITE):
    sp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()


def add_footer(slide, text=FOOTER_TEXT):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.17), Inches(13.33), Inches(0.33)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BG
    bar.line.fill.background()
    add_tb(slide, text, 0.2, 7.19, 12.9, 0.28,
           size=8, color=LIGHT_GRAY)


def add_tb(slide, text, left, top, width, height,
           size=13, bold=False, color=GRAY_TEXT,
           align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def slide_title(slide, title, subtitle=None):
    """Standard content slide title: large dark navy + gray subtitle."""
    add_tb(slide, title, 0.45, 0.3, 12.4, 0.75,
           size=34, bold=True, color=DARK_NAVY)
    if subtitle:
        add_tb(slide, subtitle, 0.45, 1.05, 12.4, 0.4,
               size=13, color=GRAY_TEXT)


def divider_line(slide, top=1.5):
    bar = slide.shapes.add_shape(
        1, Inches(0.45), Inches(top), Inches(12.4), Inches(0.03)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xDD, 0xE3, 0xEE)
    bar.line.fill.background()


def section_label(slide, text, left=0.45, top=1.6, color=MID_BLUE):
    add_tb(slide, text, left, top, 12.0, 0.38,
           size=12, bold=True, color=color)


def card(slide, left, top, width, height, bg=OFF_WHITE, border=None):
    sp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg
    if border:
        sp.line.color.rgb = border
    else:
        sp.line.fill.background()
    return sp


def stat_card(slide, value, label, left, top, w=2.8, h=1.4):
    card(slide, left, top, w, h, bg=CARD_DARK)
    add_tb(slide, value, left+0.1, top+0.08, w-0.2, 0.7,
           size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(slide, label, left+0.1, top+0.82, w-0.2, 0.5,
           size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def numbered_finding(slide, num, heading, body, left, top, w=12.0):
    add_tb(slide, num, left, top, 0.7, 0.55,
           size=30, bold=True, color=RED)
    add_tb(slide, heading, left+0.75, top+0.02, w-0.75, 0.38,
           size=13, bold=True, color=DARK_NAVY)
    add_tb(slide, body, left+0.75, top+0.42, w-0.75, 0.5,
           size=12, color=GRAY_TEXT)


def bullets(slide, items, left, top, width, height,
            size=12, color=GRAY_TEXT, spacing=Pt(5)):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing


def icon_card(slide, icon, heading, body, left, top, w=3.9, h=2.3):
    card(slide, left, top, w, h, bg=OFF_WHITE,
         border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(slide, icon,    left+0.15, top+0.12, 0.5,  0.5,  size=28)
    add_tb(slide, heading, left+0.15, top+0.65, w-0.3, 0.38,
           size=13, bold=True, color=MID_BLUE)
    add_tb(slide, body,    left+0.15, top+1.05, w-0.3, 1.1,
           size=11, color=GRAY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
#  PAPER 1 — Narrative Centrality & Competitive Outcomes
# ══════════════════════════════════════════════════════════════════════════════

def build_paper1():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    # left accent bar
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = MID_BLUE; acc.line.fill.background()

    add_tb(sl, "NARRATIVE CENTRALITY", 0.55, 0.5, 12.0, 0.5,
           size=13, bold=True, color=MID_BLUE)
    add_tb(sl, "& Competitive Outcomes\nin Major League Soccer",
           0.55, 0.95, 12.0, 1.6, size=48, bold=True, color=DARK_NAVY)
    add_tb(sl, "Why Press Coverage Predicts Playoffs but Not Points",
           0.55, 2.75, 12.0, 0.55, size=20, color=GRAY_TEXT)

    # stat boxes
    stat_card(sl, "7,236",  "Press Articles",      left=0.55, top=3.7, w=2.7)
    stat_card(sl, "186",    "Club-Season Obs.",     left=3.45, top=3.7, w=2.7)
    stat_card(sl, "29",     "MLS Clubs",            left=6.35, top=3.7, w=2.7)
    stat_card(sl, "2018–24","7 Seasons of Data",    left=9.25, top=3.7, w=2.7)

    add_tb(sl, "Shoghanian & Tang  ·  Target: Journal of Sports Analytics / ESMQ",
           0.55, 5.35, 12.0, 0.4, size=11, color=GRAY_TEXT)
    add_footer(sl)

    # ── Slide 2: The Problem ──────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "The Question",
                "Does press coverage tell us which MLS clubs will compete — and how?")
    divider_line(sl, top=1.48)
    add_footer(sl)

    icon_card(sl, "⚽", "EPL Has an Answer",
              "Our prior EPL research found narrative centrality predicts points per game — even after controlling for transfer spending.",
              left=0.45, top=1.65, w=3.9)
    icon_card(sl, "🏗️", "MLS is Different",
              "Hard salary cap, single-entity structure, playoff format. Financial resources are deliberately compressed by design.",
              left=4.55, top=1.65, w=3.9)
    icon_card(sl, "❓", "So What Changes?",
              "If spending is capped and competition is compressed, does press narrative still matter — and for which outcomes?",
              left=8.65, top=1.65, w=4.2)

    add_tb(sl, "THREE HYPOTHESES", 0.45, 4.15, 12.0, 0.35,
           size=11, bold=True, color=RED)
    bullets(sl, [
        "H1  Press narrative centrality does NOT predict next-season points in MLS",
        "H2  Press narrative centrality DOES predict same-season playoff qualification",
        "H3  The difference is explained by the salary cap compressing performance variance",
    ], left=0.45, top=4.55, width=12.4, height=2.0, size=13, color=DARK_NAVY)

    # ── Slide 3: Data ─────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "What We Built",
                "An end-to-end NLP pipeline from 7,236 articles to a 186-observation analysis panel")
    divider_line(sl, top=1.48)
    add_footer(sl)

    steps = [
        ("PRESS ARTICLES", "7,236 articles\n2018–2024\nMLS coverage"),
        ("NER PIPELINE",   "spaCy entity\nextraction\nclub + player mentions"),
        ("CO-OCC NETWORK", "Annual club\nco-occurrence\ngraph (networkx)"),
        ("CENTRALITY",     "PageRank,\ndegree, strength\nper club per year"),
        ("PANEL",          "186 obs\n29 clubs × 7 yrs\n+ performance data"),
    ]
    arrow_color = RED
    for i, (label, detail) in enumerate(steps):
        x = 0.3 + i * 2.55
        card(sl, x, 1.65, 2.3, 2.8, bg=CARD_DARK)
        add_tb(sl, label,  x+0.1, 1.75, 2.1, 0.45, size=11, bold=True,
               color=RED, align=PP_ALIGN.CENTER)
        add_tb(sl, detail, x+0.1, 2.22, 2.1, 1.8,  size=11, color=LIGHT_GRAY,
               align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_tb(sl, "→", x+2.3, 2.7, 0.25, 0.4, size=18, bold=True,
                   color=arrow_color, align=PP_ALIGN.CENTER)

    section_label(sl, "KEY MEASURE", top=4.65)
    add_tb(sl,
        "press_strength_z  =  within-year z-score of total co-occurrence edge weight\n"
        "Captures where a club sits in the discourse ecosystem relative to its peers that season — not just how often it's mentioned.",
        0.45, 5.05, 12.4, 1.0, size=13, color=DARK_NAVY)

    # ── Slide 4: Models ───────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Two Models",
                "One for continuous performance (points), one for the binary threshold (playoffs)")
    divider_line(sl, top=1.48)
    add_footer(sl)

    card(sl, 0.45, 1.65, 5.9, 4.9, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "MODEL A", 0.65, 1.75, 5.5, 0.4, size=12, bold=True, color=MID_BLUE)
    add_tb(sl, "Next-Season Points (OLS)", 0.65, 2.15, 5.5, 0.45,
           size=18, bold=True, color=DARK_NAVY)
    bullets(sl, [
        "points(t+1) = narrative + controls + Club FE + Year FE",
        "",
        "3 specs: naive → add fixed effects → add xPoints + payroll",
        "HC1 robust standard errors throughout",
        "LOYO cross-validation for temporal robustness",
    ], left=0.65, top=2.65, width=5.5, height=3.0, size=12)

    card(sl, 6.95, 1.65, 5.9, 4.9, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "MODEL B", 7.15, 1.75, 5.5, 0.4, size=12, bold=True, color=MID_BLUE)
    add_tb(sl, "Playoff Qualification (Logit)", 7.15, 2.15, 5.5, 0.45,
           size=18, bold=True, color=DARK_NAVY)
    bullets(sl, [
        "P(playoff) = narrative + xPoints + Year FE",
        "",
        "Club FE excluded (incidental parameters problem)",
        "Report marginal effects at the mean",
        "LOYO: accuracy + AUC per held-out season",
    ], left=7.15, top=2.65, width=5.5, height=3.0, size=12)

    # ── Slide 5: Finding 1 — Points null ─────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Finding 1: Narrative Does Not Predict Points",
                "Press centrality carries no signal for next-season points — across any specification")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "0.140", "Naive R² (Model A1)", left=0.45, top=1.65, w=3.0)
    stat_card(sl, "0.661", "R² with Club + Year FE", left=3.6,  top=1.65, w=3.0)
    stat_card(sl, "−1.006","LOYO Mean R²",           left=6.75, top=1.65, w=3.0)
    stat_card(sl, "n.s.",  "press_strength_z (all specs)", left=9.9, top=1.65, w=2.9)

    section_label(sl, "WHAT THIS MEANS", top=3.3)
    bullets(sl, [
        "Adding club and year fixed effects (A2) captures most of the variance — club identity matters more than narrative",
        "Once we control for xPoints and payroll (A3), squad quality and spending explain the rest",
        "press_strength_z is non-significant in every single specification — this null result is robust",
        "LOYO R² of −1.006 confirms the model doesn't generalize out-of-sample either",
        "Conclusion: in a salary-capped league, narrative alone cannot reliably predict a club's points total",
    ], left=0.45, top=3.72, width=12.4, height=3.0, size=13)

    # ── Slide 6: Finding 2 — Playoffs ────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Finding 2: Narrative Does Predict Playoff Qualification",
                "A 1 SD increase in press centrality raises the probability of making the playoffs by 82 points")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "4.82*",  "press_strength_z β\n(B3, p=.048)",  left=0.45, top=1.65, w=3.0)
    stat_card(sl, "82pp",   "Marginal Effect\nat the mean",       left=3.6,  top=1.65, w=3.0)
    stat_card(sl, "0.747",  "LOYO Mean AUC\n(7 seasons)",         left=6.75, top=1.65, w=3.0)
    stat_card(sl, "0.779",  "Mean AUC\nexcluding 2020",           left=9.9,  top=1.65, w=2.9)

    section_label(sl, "ACROSS ALL THREE LOGIT SPECS", top=3.3)
    bullets(sl, [
        "B1 (naive):        β=4.26, p=.026 — significant even without any controls",
        "B2 (+ year FE):    β=5.75, p=.009 — effect strengthens when we account for league-wide year shifts",
        "B3 (+ xPoints):    β=4.82, p=.048 — survives controlling for actual squad quality",
        "xPoints is the strongest predictor (β=0.132, p<.001) — narrative adds on top of talent, not instead of it",
        "LOYO AUC range: 0.564 in 2020 (COVID disruption) to 0.917 in the best fold",
    ], left=0.45, top=3.72, width=12.4, height=3.0, size=13)

    # ── Slide 7: Why the gap? ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Why Playoffs but Not Points?",
                "The MLS salary cap changes what narrative can and cannot do")
    divider_line(sl, top=1.48)
    add_footer(sl)

    card(sl, 0.45, 1.65, 5.9, 4.9, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "THRESHOLD MECHANISM", 0.65, 1.78, 5.5, 0.38,
           size=11, bold=True, color=RED)
    add_tb(sl, "The playoffs are a binary cliff", 0.65, 2.2, 5.5, 0.45,
           size=17, bold=True, color=DARK_NAVY)
    bullets(sl, [
        "~55% of teams make the playoffs — a binary pass/fail",
        "Narrative centrality helps attract DPs, retain fans, maintain partnerships",
        "Those marginal advantages are just enough to clear the threshold...",
        "...but not enough to push from mid-table to top in a cap world",
    ], left=0.65, top=2.7, width=5.5, height=2.8, size=12)

    card(sl, 6.95, 1.65, 5.9, 4.9, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "VARIANCE COMPRESSION", 7.15, 1.78, 5.5, 0.38,
           size=11, bold=True, color=RED)
    add_tb(sl, "The cap flattens the playing field", 7.15, 2.2, 5.5, 0.45,
           size=17, bold=True, color=DARK_NAVY)
    bullets(sl, [
        "MLS points SD ≈ 11  vs  EPL ≈ 20+ pts",
        "Top/bottom payroll ratio: 3× in MLS vs 20× in EPL",
        "In EPL, narrative-driven resources compound into a full-season edge",
        "In MLS, that same edge only shows up at the binary playoff cutoff",
    ], left=7.15, top=2.7, width=5.5, height=2.8, size=12)

    # ── Slide 8: Press vs Reddit ──────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Two Parallel Narratives",
                "Press and Reddit track very different things — and that gap tells a story")
    divider_line(sl, top=1.48)
    add_footer(sl)

    quad_data = [
        (MID_BLUE, "Institutional Legacy",
         "Toronto FC  ·  LA Galaxy  ·  CF Montreal",
         "High press, low Reddit. Historical success sustains coverage long after competitive decline."),
        (RED, "Grassroots Underdogs",
         "Philadelphia Union  ·  Columbus Crew",
         "Fan discourse outpaces press. Active communities reflect real competitive quality."),
        (RGBColor(0x22, 0x7A, 0x3A), "Consensus Stars",
         "Inter Miami CF  ·  Seattle Sounders",
         "High in both channels. Narrative and performance are aligned."),
        (RGBColor(0x66, 0x66, 0x88), "Low Profile",
         "Nashville  ·  San Jose  ·  St. Louis  ·  Colorado",
         "Below median in both. Expansion clubs and small markets not yet established."),
    ]
    for i, (col, title, clubs, desc) in enumerate(quad_data):
        x = 0.45 + i * 3.22
        card(sl, x, 1.65, 3.0, 5.0, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
        bar = sl.shapes.add_shape(1, Inches(x), Inches(1.65), Inches(3.0), Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
        add_tb(sl, title, x+0.1, 1.85, 2.8, 0.45, size=13, bold=True, color=col)
        add_tb(sl, clubs, x+0.1, 2.35, 2.8, 0.5,  size=11, bold=True, color=DARK_NAVY)
        add_tb(sl, desc,  x+0.1, 2.9,  2.8, 2.5,  size=11, color=GRAY_TEXT)

    # ── Slide 9: Toronto FC case ──────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Case Study: Toronto FC",
                "The clearest example of institutional inertia — press stuck in 2017, Reddit in 2024")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "6.0",   "Press Rank (avg)\n2018–2024",      left=0.45, top=1.65, w=2.9)
    stat_card(sl, "19.4",  "Reddit Rank (avg)\n2018–2024",     left=3.5,  top=1.65, w=2.9)
    stat_card(sl, "Bot. 3","EC Finish\nevery yr 2020–2024",    left=6.55, top=1.65, w=2.9)
    stat_card(sl, "+1→+21","Rank Gap\n2018 → 2022",            left=9.6,  top=1.65, w=2.9)

    section_label(sl, "WHAT HAPPENED", top=3.3)
    bullets(sl, [
        "Toronto won the MLS Treble in 2017 — it was the biggest story in league history at the time",
        "That narrative capital kept them in the top 6 of press coverage for the entire 7-year panel",
        "Meanwhile they finished bottom 3 in the Eastern Conference every single season from 2020 to 2024",
        "Reddit updated in real time — fans knew the team was bad. Press took years to catch up.",
        "The gap widened every year: +1 in 2018, growing to +21 by 2022 and holding there through 2023",
        "By 2024 the press rank had declined slightly to 4th — still top-5, still above a perennial also-ran",
    ], left=0.45, top=3.72, width=12.4, height=3.0, size=13)

    # ── Slide 10: Cross-league ────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "MLS vs EPL: Two Different Worlds",
                "The same narrative measure behaves differently depending on the league's institutional structure")
    divider_line(sl, top=1.48)
    add_footer(sl)

    headers = ["", "MLS (This Study)", "EPL (Shoghanian & Tang, forthcoming)"]
    rows = [
        ["Salary structure",         "Hard cap — top/bottom ≈ 3×",   "No cap — top/bottom ≈ 20×"],
        ["Narrative → Points",       "Non-significant (all specs)",   "Significant after controls"],
        ["Narrative → Playoffs",     "Significant β=4.82, p=.048",   "N/A (table-based system)"],
        ["Out-of-sample performance","LOYO AUC = 0.747",              "R² holds cross-season"],
        ["Press inertia",            "Strong (Toronto FC)",           "Present but more moderate"],
    ]
    col_w = [3.2, 4.4, 4.4]
    col_x = [0.45, 3.85, 8.45]
    rh = 0.68
    y0 = 1.65

    for ci, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        hb = sl.shapes.add_shape(1, Inches(x), Inches(y0), Inches(w), Inches(rh*0.92))
        hb.fill.solid()
        hb.fill.fore_color.rgb = CARD_DARK if ci > 0 else RGBColor(0xEE, 0xF2, 0xF8)
        hb.line.fill.background()
        add_tb(sl, hdr, x+0.1, y0+0.06, w-0.2, rh*0.8,
               size=13, bold=True,
               color=WHITE if ci > 0 else GRAY_TEXT,
               align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        y = y0 + rh * (ri + 1)
        bg = WHITE if ri % 2 == 0 else RGBColor(0xF0, 0xF3, 0xF9)
        for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            cb = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(rh*0.92))
            cb.fill.solid(); cb.fill.fore_color.rgb = bg
            cb.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEE)
            bold = (ci == 0)
            add_tb(sl, cell, x+0.1, y+0.07, w-0.2, rh*0.8,
                   size=12, bold=bold,
                   color=DARK_NAVY if ci == 0 else GRAY_TEXT)

    # ── Slide 11: Key Findings Summary ───────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Key Findings",
                "Four evidence-backed takeaways from the MLS narrative analysis")
    divider_line(sl, top=1.48)
    add_footer(sl)

    numbered_finding(sl, "01",
        "Narrative does not predict points",
        "press_strength_z is non-significant across all OLS specifications. Salary cap compresses performance variance.",
        left=0.45, top=1.65)
    numbered_finding(sl, "02",
        "Narrative does predict playoff qualification",
        "β=4.82 (p=.048), marginal effect of 82pp. LOYO mean AUC = 0.747 across seven seasons.",
        left=0.45, top=2.7)
    numbered_finding(sl, "03",
        "Press and Reddit carry different information",
        "Press tracks legacy and reputation. Reddit tracks current performance. The gap reveals club type.",
        left=0.45, top=3.75)
    numbered_finding(sl, "04",
        "The EPL contrast validates the theory",
        "Same measure predicts points in uncapped EPL but only playoff threshold in capped MLS. Structure mediates narrative.",
        left=0.45, top=4.8)

    # ── Slide 12: Bottom line ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, CARD_DARK)
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = MID_BLUE; acc.line.fill.background()

    add_tb(sl, "The Bottom Line", 0.55, 0.5, 12.0, 0.7,
           size=34, bold=True, color=WHITE)
    add_tb(sl,
        "Press narrative centrality predicts whether MLS clubs make the playoffs.\n"
        "It does not predict how many points they earn — and that difference is explained\n"
        "by the salary cap flattening the path from narrative to performance.",
        0.55, 1.4, 11.5, 1.8, size=20, color=LIGHT_GRAY)

    add_tb(sl, "For Clubs", 0.55, 3.4, 3.8, 0.4,  size=16, bold=True, color=MID_BLUE)
    add_tb(sl, "Narrative management helps clear the playoff threshold. That's the margin that matters in MLS.",
           0.55, 3.85, 3.8, 1.2, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Researchers", 4.65, 3.4, 3.8, 0.4, size=16, bold=True, color=MID_BLUE)
    add_tb(sl, "League institutional structure mediates how narrative translates to outcomes. Test across leagues.",
           4.65, 3.85, 3.8, 1.2, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Future Work", 8.75, 3.4, 3.8, 0.4, size=16, bold=True, color=MID_BLUE)
    add_tb(sl, "Apply to NWSL, MLS NEXT Pro. Test whether DP signings amplify the narrative-playoff link.",
           8.75, 3.85, 3.8, 1.2, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

    add_footer(sl)

    prs.save(OUT_DIR / "paper1_slides.pptx")
    print("Saved: presentations/paper1_slides.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PAPER 2 — Earned Narrative & Commercial Returns
# ══════════════════════════════════════════════════════════════════════════════

def build_paper2():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = RED; acc.line.fill.background()

    add_tb(sl, "EARNED NARRATIVE", 0.55, 0.5, 12.0, 0.5,
           size=13, bold=True, color=RED)
    add_tb(sl, "& Commercial Returns\nin Major League Soccer",
           0.55, 0.95, 12.0, 1.6, size=48, bold=True, color=DARK_NAVY)
    add_tb(sl, "Evidence from Press Network Centrality",
           0.55, 2.75, 12.0, 0.55, size=20, color=GRAY_TEXT)

    stat_card(sl, "+21.9%", "Valuation premium\nper 1 SD narrative",  left=0.55, top=3.7, w=2.7)
    stat_card(sl, "+18.9%", "Revenue premium\nper 1 SD narrative",    left=3.45, top=3.7, w=2.7)
    stat_card(sl, "0.629",  "Model R²\n(valuation, year FE)",         left=6.35, top=3.7, w=2.7)
    stat_card(sl, "n.s.",   "On-field points\n(not significant)",      left=9.25, top=3.7, w=2.7)

    add_tb(sl, "Shoghanian & Tang  ·  Target: Sport Marketing Quarterly",
           0.55, 5.35, 12.0, 0.4, size=11, color=GRAY_TEXT)
    add_footer(sl)

    # ── Slide 2: The Puzzle ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "The Commercial Puzzle",
                "MLS club valuations vary by 10× — but the salary cap deliberately compresses on-field performance")
    divider_line(sl, top=1.48)
    add_footer(sl)

    add_tb(sl,
        "If every team plays under the same financial constraints, why are some clubs worth ten times more than others?",
        0.45, 1.65, 12.4, 0.8, size=18, bold=True, color=DARK_NAVY)

    icon_card(sl, "📰", "Narrative Capital Hypothesis",
              "A club's structural position in the press co-occurrence network = earned media capital that signals value to sponsors, players, and investors.",
              left=0.45, top=2.65, w=3.9)
    icon_card(sl, "💵", "Commercial Outcomes We Test",
              "Club valuation (Forbes), annual revenue (Forbes), jersey sponsorship deal value, and total player payroll (MLSPA).",
              left=4.55, top=2.65, w=3.9)
    icon_card(sl, "🎯", "The Core Claim",
              "Narrative centrality predicts commercial value independently of points, attendance, and market size. Narrative is a measurable asset.",
              left=8.65, top=2.65, w=4.2)

    section_label(sl, "THREE CONTRIBUTIONS", top=5.0)
    bullets(sl, [
        "1.  First network-based NLP analysis of MLS press coverage spanning 2018–2024",
        "2.  Demonstrate narrative centrality predicts valuation, revenue, and payroll independently of performance",
        "3.  Document an asymmetry between press and Reddit — each channel carries distinct commercial signals",
    ], left=0.45, top=5.42, width=12.4, height=1.5, size=12)

    # ── Slide 3: Data ─────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Data Sources",
                "186 club-season observations across 29 clubs and 7 seasons (2018–2024)")
    divider_line(sl, top=1.48)
    add_footer(sl)

    sources = [
        ("📰 Press",       "7,236 articles\nco-occurrence NLP"),
        ("💬 Reddit",      "28 subreddits\nclub co-mentions"),
        ("💰 Valuations",  "Forbes MLS\n164 obs"),
        ("📈 Revenue",     "Forbes MLS\n164 obs"),
        ("💼 Payroll",     "MLSPA 2018–24\n(excl. 2020)"),
        ("🏷️ Sponsors",   "Jersey deals\nfrom announcements"),
        ("⚽ Performance", "Official MLS\n+ ASA xGoals"),
    ]
    for i, (label, detail) in enumerate(sources):
        r, c = divmod(i, 4)
        x = 0.45 + c * 3.22
        y = 1.65 + r * 2.4
        card(sl, x, y, 3.0, 2.15, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
        add_tb(sl, label,  x+0.15, y+0.15, 2.7, 0.5,  size=14, bold=True, color=DARK_NAVY)
        add_tb(sl, detail, x+0.15, y+0.7,  2.7, 1.25, size=12, color=GRAY_TEXT)

    # ── Slide 4: Model ────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Empirical Strategy",
                "OLS regression with year fixed effects and HC1 robust standard errors")
    divider_line(sl, top=1.48)
    add_footer(sl)

    card(sl, 0.45, 1.65, 12.4, 1.6, bg=RGBColor(0xF0, 0xF3, 0xF9),
         border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl,
        "log(Y) = α  +  β₁·press_strength_z  +  β₂·press/reddit_ratio\n"
        "             +  β₃·sentiment_gap  +  β₄·press_net_momentum\n"
        "             +  γ·[points, log-attendance, conference]  +  δ_year  +  ε",
        0.65, 1.8, 12.0, 1.35, size=15, color=DARK_NAVY)

    section_label(sl, "FOUR OUTCOME MODELS", top=3.45)
    outcomes = [
        ("C1", "log(Valuation)",  "Forbes club valuation"),
        ("C2", "log(Revenue)",    "Annual club revenue"),
        ("C3", "Sponsor Deal $",  "Jersey deal value"),
        ("C4", "log(Payroll)",    "Total player payroll"),
    ]
    for i, (code, name, desc) in enumerate(outcomes):
        x = 0.45 + i * 3.22
        card(sl, x, 3.85, 3.0, 2.8, bg=CARD_DARK)
        add_tb(sl, code, x+0.15, 3.98, 2.7, 0.5,  size=22, bold=True, color=RED,
               align=PP_ALIGN.CENTER)
        add_tb(sl, name, x+0.15, 4.5,  2.7, 0.45, size=14, bold=True, color=WHITE,
               align=PP_ALIGN.CENTER)
        add_tb(sl, desc, x+0.15, 4.98, 2.7, 1.4,  size=11, color=LIGHT_GRAY,
               align=PP_ALIGN.CENTER)

    # ── Slide 5: Main Results ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Finding 1: Narrative Drives Valuation and Revenue",
                "A 1 SD increase in press network strength is associated with 21.9% higher club valuations")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "0.198***", "press_strength_z β\nC1: log(valuation)", left=0.45, top=1.65, w=3.0)
    stat_card(sl, "+21.9%",   "Valuation premium\nper 1 SD narrative",  left=3.6,  top=1.65, w=3.0)
    stat_card(sl, "0.173***", "press_strength_z β\nC2: log(revenue)",   left=6.75, top=1.65, w=3.0)
    stat_card(sl, "n.s.",     "Points coefficient\n(not significant)",   left=9.9,  top=1.65, w=2.9)

    section_label(sl, "WHAT THE TABLE SHOWS", top=3.3)
    bullets(sl, [
        "Moving from 25th to 75th percentile in narrative centrality ≈ $83M valuation premium at the sample mean",
        "Points are NOT significant once narrative and attendance are controlled — winning doesn't drive value here",
        "Attendance (β=0.142, p=.008) and Eastern Conference (β=0.108, p=.024) are significant controls",
        "press_net_momentum is NEGATIVE (β=−0.031, p=.012) — ascending clubs are smaller and lower-valued today",
        "C3 (sponsor deals): directional positive but only marginal (p=.068) — long contracting cycles reduce sensitivity",
    ], left=0.45, top=3.72, width=12.4, height=3.0, size=13)

    # ── Slide 6: Payroll finding ──────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Finding 2: Narrative Shapes Investment Decisions",
                "Clubs with more press coverage AND more fan engagement both attract higher payrolls")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "0.163***", "press_strength_z β\nC4: log(payroll)",      left=0.45, top=1.65, w=3.9)
    stat_card(sl, "−0.079**", "press/reddit ratio β\n(negative = fan-driven)", left=4.55, top=1.65, w=3.9)
    stat_card(sl, "0.538",    "Model R²\n(payroll model)",                  left=8.65, top=1.65, w=3.9)

    section_label(sl, "THE ASYMMETRY", top=3.3)
    add_tb(sl,
        "press_strength_z is positive: more press coverage → higher payroll\n"
        "press_to_reddit_ratio is negative: more fan-driven narrative (Reddit-heavy clubs) → also higher payroll",
        0.45, 3.72, 12.4, 1.0, size=14, bold=True, color=DARK_NAVY)

    section_label(sl, "WHAT THIS MEANS", top=4.85)
    bullets(sl, [
        "Organic fan engagement signals market depth to ownership — independent of what the press is saying",
        "Press coverage can reflect institutional inertia (Toronto's 2017 treble still generating headlines in 2024)",
        "Reddit is harder to manufacture — it reflects real current fan enthusiasm",
        "Implication: narrative REACH (press) drives external valuations; narrative AUTHENTICITY (fans) drives investment",
    ], left=0.45, top=5.27, width=12.4, height=1.8, size=13)

    # ── Slide 7: LOYO ─────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Robustness — and an Honest Limitation",
                "The model works well pre-2022 but degrades sharply after — and that degradation tells us something")
    divider_line(sl, top=1.48)
    add_footer(sl)

    loyo = [("2018","0.888",True),("2019","0.916",True),("2020","0.343",False),
            ("2021","0.203",False),("2022","−0.774",False),("2023","−2.116",False),("2024","−4.273",False)]
    for i, (yr, r2, good) in enumerate(loyo):
        bg = RGBColor(0x15, 0x50, 0x2E) if good else RGBColor(0x6B, 0x0F, 0x1A)
        x = 0.45 + i * 1.78
        card(sl, x, 1.65, 1.65, 1.3, bg=bg)
        add_tb(sl, r2,  x+0.1, 1.72, 1.45, 0.65, size=20, bold=True, color=WHITE,
               align=PP_ALIGN.CENTER)
        add_tb(sl, yr,  x+0.1, 2.35, 1.45, 0.45, size=10, color=LIGHT_GRAY,
               align=PP_ALIGN.CENTER)

    section_label(sl, "WHY THE MODEL BREAKS AFTER 2022", top=3.15)
    bullets(sl, [
        "2022+: MLS expansion fees hit $200M+, new broadcast rights, Messi arrives at Inter Miami in 2023",
        "League-wide valuation inflation means every club's value jumped — compressing the cross-sectional variance the model uses",
        "This is a limitation WITH informational content: it signals a commercial regime change in MLS post-2021",
        "The core finding (narrative → valuation) is robustly significant within-sample across all seven years",
        "Future models should include a post-2021 structural break indicator or re-estimate from 2022 forward separately",
    ], left=0.45, top=3.57, width=12.4, height=3.0, size=13)

    # ── Slide 8: Key Findings ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Key Findings",
                "Four commercial insights from seven seasons of MLS press network data")
    divider_line(sl, top=1.48)
    add_footer(sl)

    numbered_finding(sl, "01",
        "Press narrative centrality predicts valuation and revenue",
        "+21.9% valuation and +18.9% revenue per 1 SD increase, independent of on-field performance.",
        left=0.45, top=1.65)
    numbered_finding(sl, "02",
        "Points don't drive commercial value — narrative does",
        "Points coefficient is non-significant once narrative and attendance are controlled.",
        left=0.45, top=2.7)
    numbered_finding(sl, "03",
        "Fan engagement (Reddit) carries independent investment signals",
        "Clubs with more fan-driven narrative attract higher payrolls — organic engagement signals market depth.",
        left=0.45, top=3.75)
    numbered_finding(sl, "04",
        "The model breaks post-2021 — and that's informative",
        "Messi, expansion fees, and broadcast deals created a commercial regime change. Historical relationships shifted.",
        left=0.45, top=4.8)

    # ── Slide 9: Bottom line ──────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, CARD_DARK)
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = RED; acc.line.fill.background()

    add_tb(sl, "The Bottom Line", 0.55, 0.5, 12.0, 0.7,
           size=34, bold=True, color=WHITE)
    add_tb(sl,
        "In a salary-capped league where winning is structurally compressed,\n"
        "narrative management is not just reputational — it's financially material.",
        0.55, 1.4, 11.5, 1.5, size=20, color=LIGHT_GRAY)

    add_tb(sl, "For Club Executives", 0.55, 3.2, 3.8, 0.4,  size=16, bold=True, color=RED)
    add_tb(sl, "Sustained press presence — not just winning — constitutes a strategic asset worth managing deliberately.",
           0.55, 3.65, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Sport Marketers", 4.65, 3.2, 3.8, 0.4, size=16, bold=True, color=RED)
    add_tb(sl, "Narrative reach and narrative authenticity matter differently. Press drives valuation; fans drive investment.",
           4.65, 3.65, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Future Research", 8.75, 3.2, 3.8, 0.4, size=16, bold=True, color=RED)
    add_tb(sl, "Replicate in NFL, NBA, NWSL. Test whether press narrative mediates fan engagement or operates independently.",
           8.75, 3.65, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

    add_footer(sl)

    prs.save(OUT_DIR / "paper2_slides.pptx")
    print("Saved: presentations/paper2_slides.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  RESEARCH NOTE — Press-Reddit Divergence
# ══════════════════════════════════════════════════════════════════════════════

def build_research_note():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = RGBColor(0x22, 0x7A, 0x3A)
    acc.line.fill.background()

    add_tb(sl, "WHEN MEDIA AND FANS DISAGREE", 0.55, 0.5, 12.0, 0.5,
           size=13, bold=True, color=RGBColor(0x22, 0x7A, 0x3A))
    add_tb(sl, "Press-Reddit Narrative\nDivergence in MLS",
           0.55, 0.95, 12.0, 1.7, size=46, bold=True, color=DARK_NAVY)
    add_tb(sl, "A Research Note on Two Parallel Ecosystems",
           0.55, 2.85, 12.0, 0.55, size=20, color=GRAY_TEXT)

    stat_card(sl, "186",     "Club-season obs.\n2018–2024",        left=0.55, top=3.7, w=2.7)
    stat_card(sl, "29",      "MLS clubs\nclassified",              left=3.45, top=3.7, w=2.7)
    stat_card(sl, "4",       "Club type\nquadrants",               left=6.35, top=3.7, w=2.7)
    stat_card(sl, "r=−0.012","Gap vs valuation\n(not predictive)", left=9.25, top=3.7, w=2.7)

    add_tb(sl, "Shoghanian & Tang  ·  Target: International Journal of Sport Communication",
           0.55, 5.35, 12.0, 0.4, size=11, color=GRAY_TEXT)
    add_footer(sl)

    # ── Slide 2: Research Questions ───────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Three Questions",
                "What happens when institutional media and fan communities tell different stories?")
    divider_line(sl, top=1.48)
    add_footer(sl)

    questions = [
        ("RQ1", "Do press and fan narratives\nsystematically diverge?",
         "Which MLS clubs show the largest gaps between what the press covers and what fans talk about?"),
        ("RQ2", "Does one channel lead\nthe other over time?",
         "Granger causality test — does press sentiment predict future Reddit sentiment, or vice versa?"),
        ("RQ3", "Does divergence predict\ncommercial outcomes?",
         "If press and fans disagree about a club, does that gap forecast future valuation changes?"),
    ]
    for i, (rq, q, sub) in enumerate(questions):
        y = 1.65 + i * 1.72
        card(sl, 0.45, y, 12.4, 1.55, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
        add_tb(sl, rq,  0.6,  y+0.12, 0.85, 1.2, size=24, bold=True, color=RED)
        add_tb(sl, q,   1.6,  y+0.12, 4.5,  0.95, size=15, bold=True, color=DARK_NAVY)
        add_tb(sl, sub, 6.3,  y+0.22, 6.3,  1.0,  size=12, color=GRAY_TEXT)

    # ── Slide 3: Method ───────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "How We Measure the Gap",
                "Two parallel co-occurrence networks, one for press and one for Reddit, built from the same clubs")
    divider_line(sl, top=1.48)
    add_footer(sl)

    card(sl, 0.45, 1.65, 5.9, 2.5, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "📰 PRESS NETWORK", 0.65, 1.78, 5.5, 0.4, size=12, bold=True, color=MID_BLUE)
    bullets(sl, [
        "7,236 articles → annual club co-occurrence graphs",
        "PageRank centrality per club per season",
        "Within-season rank assigned (rank 1 = most central)",
    ], left=0.65, top=2.22, width=5.5, height=1.6, size=12)

    card(sl, 6.95, 1.65, 5.9, 2.5, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "💬 REDDIT NETWORK", 7.15, 1.78, 5.5, 0.4, size=12, bold=True, color=MID_BLUE)
    bullets(sl, [
        "r/MLS + r/ussoccer + 26 club-specific subreddits",
        "Club co-mention networks → Reddit PageRank ranks",
        "Same within-season rank structure",
    ], left=7.15, top=2.22, width=5.5, height=1.6, size=12)

    section_label(sl, "THE DIVERGENCE MEASURE", top=4.35)
    add_tb(sl,
        "Press-Reddit Rank Gap  =  press rank  −  Reddit rank",
        0.45, 4.78, 12.4, 0.5, size=16, bold=True, color=DARK_NAVY)
    bullets(sl, [
        "Positive gap → club is more prominent on Reddit than in press  (Grassroots Underdog)",
        "Negative gap → club is more prominent in press than on Reddit  (Institutional Legacy)",
    ], left=0.45, top=5.32, width=12.4, height=1.2, size=13)

    # ── Slide 4: Four Club Types ──────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Four Club Types",
                "A median split on average press rank vs Reddit rank classifies all 29 clubs")
    divider_line(sl, top=1.48)
    add_footer(sl)

    quads = [
        (MID_BLUE,  "INSTITUTIONAL LEGACY",  "High press · Low Reddit",
         "Toronto FC\nLA Galaxy\nCF Montreal",
         "Historical capital keeps these clubs prominent in press long after their competitive peak."),
        (RED,       "GRASSROOTS UNDERDOGS",  "High Reddit · Lower press",
         "Philadelphia Union\nColumbus Crew",
         "Fan communities are active and sophisticated — engagement exceeds the media footprint."),
        (RGBColor(0x22,0x7A,0x3A), "CONSENSUS STARS", "High in both channels",
         "Inter Miami CF\nSeattle Sounders",
         "Narrative and performance are aligned. Both channels agree these clubs matter."),
        (GRAY_TEXT, "LOW PROFILE",           "Low in both channels",
         "Nashville · San Jose\nSt. Louis · Colorado",
         "Expansion clubs and small markets that haven't established a narrative footprint yet."),
    ]
    for i, (col, title, sub, clubs, desc) in enumerate(quads):
        x = 0.45 + i * 3.22
        card(sl, x, 1.65, 3.0, 5.1, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
        bar2 = sl.shapes.add_shape(1, Inches(x), Inches(1.65), Inches(3.0), Inches(0.1))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = col; bar2.line.fill.background()
        add_tb(sl, title,  x+0.12, 1.83, 2.75, 0.45, size=12, bold=True, color=col)
        add_tb(sl, sub,    x+0.12, 2.3,  2.75, 0.38, size=11, color=GRAY_TEXT)
        add_tb(sl, clubs,  x+0.12, 2.75, 2.75, 0.9,  size=13, bold=True, color=DARK_NAVY)
        add_tb(sl, desc,   x+0.12, 3.72, 2.75, 2.7,  size=11, color=GRAY_TEXT)

    # ── Slide 5: Case Studies ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "Five Case Studies",
                "Each club illustrates a different version of how press and fan narratives can diverge")
    divider_line(sl, top=1.48)
    add_footer(sl)

    cases = [
        ("Toronto FC",       "Institutional Legacy",
         "Press rank: 6.0  ·  Reddit: 19.4\nGap: +1 (2018) → +21 (2022)\nBottom-3 EC every year 2020–2024\nStill ranked top-4 in press in 2024"),
        ("CF Montreal",      "Institutional Legacy",
         "Press rank: 3.9  ·  Reddit: 12.0\nNote: r/cfmontreal not in dataset\nReddit side likely underestimated\nConference finalist in 2021"),
        ("Philadelphia Union","Grassroots Underdog",
         "Reddit rank: 6.0  ·  Press rank: 9\nMLS Cup finalist 2022\nTop-5 xPoints multiple seasons\nFan community outpaces coverage"),
        ("Charlotte FC",     "Expansion Over-hype",
         "Press rank 22 at launch · Reddit rank 8\nNew stadium + DP signings = press\nFan community depth takes years\nGap closes as novelty fades"),
        ("Inter Miami / Seattle","Consensus Stars",
         "High in both channels\nSeattle: 3 MLS Cup appearances\nInter Miami: Messi arrival 2023\nStrong press + fan alignment"),
    ]
    for i, (name, label, detail) in enumerate(cases):
        r, c = divmod(i, 3)
        x = 0.45 + c * 4.3
        y = 1.65 + r * 2.75
        card(sl, x, y, 4.1, 2.55, bg=OFF_WHITE, border=RGBColor(0xDD, 0xE3, 0xEE))
        add_tb(sl, name,   x+0.15, y+0.12, 3.8, 0.45, size=14, bold=True, color=DARK_NAVY)
        add_tb(sl, label,  x+0.15, y+0.55, 3.8, 0.35, size=11, bold=True, color=MID_BLUE)
        add_tb(sl, detail, x+0.15, y+0.92, 3.8, 1.5,  size=11, color=GRAY_TEXT)

    # ── Slide 6: Granger ─────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "RQ2: Does One Channel Lead the Other?",
                "Granger causality test on the 7-year aggregate sentiment series")
    divider_line(sl, top=1.48)
    add_footer(sl)

    stat_card(sl, "0.461–0.576", "Press VADER\nsentiment range",  left=0.45, top=1.65, w=3.0)
    stat_card(sl, "0.239–0.305", "Reddit VADER\nsentiment range", left=3.6,  top=1.65, w=3.0)
    stat_card(sl, "p=.718",      "Press→Reddit\nGranger test",    left=6.75, top=1.65, w=3.0)
    stat_card(sl, "p=.603",      "Reddit→Press\nGranger test",    left=9.9,  top=1.65, w=2.9)

    section_label(sl, "WHAT THIS TELLS US", top=3.3)
    bullets(sl, [
        "Press sentiment is consistently higher than Reddit — journalists are more positive and promotional than fans",
        "Both series are non-stationary (ADF: press p=.125, Reddit p=.256) — they move together but don't trend",
        "Neither Granger test is significant — press doesn't lead Reddit, and Reddit doesn't lead press",
        "Most likely explanation: both respond to the same events (COVID, Messi, expansion) simultaneously",
        "Critical caveat: N=7 annual observations is too small for reliable Granger testing",
        "The monthly series (N=84) is available in this dataset — that's where future work should go",
    ], left=0.45, top=3.72, width=12.4, height=3.0, size=13)

    # ── Slide 7: RQ3 + Conclusions ────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_title(sl, "RQ3 + Key Takeaways",
                "The press-Reddit gap characterizes club types — but does not forecast commercial outcomes")
    divider_line(sl, top=1.48)
    add_footer(sl)

    card(sl, 0.45, 1.65, 12.4, 1.4, bg=RGBColor(0xF0, 0xF3, 0xF9),
         border=RGBColor(0xDD, 0xE3, 0xEE))
    add_tb(sl, "RQ3 Answer:", 0.65, 1.75, 3.0, 0.4, size=14, bold=True, color=DARK_NAVY)
    add_tb(sl,
        "r = −0.012 between rank gap and next-year valuation change   →   No predictive relationship",
        0.65, 2.18, 11.8, 0.5, size=15, bold=True, color=RED)

    section_label(sl, "THE GAP IS A DIAGNOSTIC TOOL, NOT A PREDICTIVE ONE", top=3.25)
    add_tb(sl, "It tells you the type of club you're dealing with. It does not tell you where the club is going commercially.",
           0.45, 3.68, 12.4, 0.55, size=13, color=DARK_NAVY)

    section_label(sl, "TAKEAWAYS", top=4.35)
    numbered_finding(sl, "01",
        "Press and fan narratives carry complementary — not redundant — information",
        "Each channel reflects a different reality. Combining them gives a fuller picture than either alone.",
        left=0.45, top=4.8)
    numbered_finding(sl, "02",
        "Institutional inertia is real and slow to correct",
        "Toronto FC: top-6 press coverage for 5 years of non-playoff finishes. The market corrects, but slowly.",
        left=0.45, top=5.85)

    # ── Slide 8: Bottom line ──────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, CARD_DARK)
    acc = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    acc.fill.solid(); acc.fill.fore_color.rgb = RGBColor(0x22, 0x7A, 0x3A)
    acc.line.fill.background()

    add_tb(sl, "The Bottom Line", 0.55, 0.5, 12.0, 0.7,
           size=34, bold=True, color=WHITE)
    add_tb(sl,
        "Press and Reddit in MLS are two different information systems.\n"
        "Press tracks legacy and institutional reputation. Reddit tracks current reality.\n"
        "The gap between them tells you what kind of club you're looking at.",
        0.55, 1.4, 11.5, 1.8, size=18, color=LIGHT_GRAY)

    add_tb(sl, "For Practitioners", 0.55, 3.4, 3.8, 0.4,  size=16, bold=True,
           color=RGBColor(0x44, 0xAA, 0x66))
    add_tb(sl, "Press narrative is manageable via media relations. Reddit is a diagnostic — it tells you what fans actually think.",
           0.55, 3.85, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Researchers", 4.65, 3.4, 3.8, 0.4, size=16, bold=True,
           color=RGBColor(0x44, 0xAA, 0x66))
    add_tb(sl, "Network-based divergence measures offer a structural way to compare institutional and fan discourse at scale.",
           4.65, 3.85, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_tb(sl, "For Future Work", 8.75, 3.4, 3.8, 0.4, size=16, bold=True,
           color=RGBColor(0x44, 0xAA, 0x66))
    add_tb(sl, "Use the monthly series (N=84) for Granger testing. Include r/cfmontreal. Apply to NWSL and PWHL.",
           8.75, 3.85, 3.8, 1.4, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

    add_footer(sl)

    prs.save(OUT_DIR / "research_note_slides.pptx")
    print("Saved: presentations/research_note_slides.pptx")


if __name__ == "__main__":
    build_paper1()
    build_paper2()
    build_research_note()
    print("\nAll three decks generated in presentations/")
